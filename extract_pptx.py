from pptx import Presentation
import json

prs = Presentation('IEEE PRESIDENCY UNIVERISTY STUDENT BRANCH CHAPTERS.pptx')
slides = []
for i, slide in enumerate(prs.slides, start=1):
    text = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text.append(shape.text)
    slides.append({'slide': i, 'text': '\n'.join(text)})
print(json.dumps(slides, indent=2, ensure_ascii=False))
