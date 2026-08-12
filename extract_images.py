from pptx import Presentation
from pathlib import Path

prs = Presentation('IEEE PRESIDENCY UNIVERISTY STUDENT BRANCH CHAPTERS.pptx')
output_dir = Path('pptx_images')
output_dir.mkdir(exist_ok=True)
count = 0
for i, slide in enumerate(prs.slides, start=1):
    for shape in slide.shapes:
        if shape.shape_type == 13:  # picture
            image = shape.image
            ext = image.ext
            path = output_dir / f'slide{i}_img{count+1}.{ext}'
            path.write_bytes(image.blob)
            count += 1
print(f'saved {count} images to {output_dir}')
